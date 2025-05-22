import argparse
import sys
import os
from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import google.generativeai as genai
import json
import re
from colorama import Fore, Style, init

# Initialize colorama
init(autoreset=True)

# Language mapping for style adaptation
LANGUAGE_MAPPINGS = {
    "en": "English",
    "ja": "Japanese",
    "es": "Spanish",
    "fr": "French",
    "de": "German",
    "zh": "Chinese",
    "ko": "Korean",
    "ru": "Russian",
    "ar": "Arabic",
    "pt": "Portuguese",
    "it": "Italian",
    # Add more languages as needed
}

# Translation style prompts
STYLE_PROMPTS = {
    "business": {
        "ja": "Keep the translation polite, collaborative, and aligned with Japanese business norms with proper Japanese business terminology..",
        "es": "Maintain a formal and professional tone suitable for Spanish business contexts.",
        "fr": "Use formal language with appropriate business terminology for French corporate environments.",
        "de": "Ensure formality with proper German business etiquette and terminology.",
        "zh": "Adapt to Chinese business communication norms with appropriate formal expressions.",
        "ko": "Follow Korean business hierarchy and formality conventions.",
        "ru": "Use formal Russian suitable for professional business communication.",
        "ar": "Employ formal Arabic with appropriate business terminology.",
        "pt": "Maintain formal Portuguese suitable for professional settings.",
        "it": "Use formal Italian appropriate for business communication.",
        "en": "Maintain professional business English with appropriate terminology.",
        # Default for other languages
        "default": "Maintain a professional, formal tone appropriate for business contexts."
    },
    "casual": {
        "ja": "Use natural, conversational Japanese while maintaining appropriate politeness.",
        "es": "Use friendly, conversational Spanish appropriate for general audiences.",
        "fr": "Use everyday French with a warm, approachable tone.",
        "de": "Employ conversational German with a friendly tone.",
        "zh": "Use natural, everyday Chinese expressions with appropriate politeness.",
        "ko": "Use conversational Korean while maintaining appropriate respect levels.",
        "ru": "Use friendly, everyday Russian for general communication.",
        "ar": "Use conversational Arabic with appropriate informality.",
        "pt": "Use natural, everyday Portuguese expressions.",
        "it": "Use casual Italian with a friendly tone.",
        "en": "Use natural, conversational English with a friendly tone.",
        # Default for other languages
        "default": "Use a casual, friendly tone while maintaining cultural appropriateness."
    },
    "technical": {
        "ja": "Prioritize technical accuracy in Japanese with appropriate terminology.",
        "es": "Use precise Spanish technical terminology and maintain factual accuracy.",
        "fr": "Employ specific French technical terms with precision and clarity.",
        "de": "Use accurate German technical terminology with precise language.",
        "zh": "Focus on technical precision in Chinese with appropriate specialized vocabulary.",
        "ko": "Use Korean technical terminology with precision and clarity.",
        "ru": "Employ specialized Russian technical vocabulary with accuracy.",
        "ar": "Use precise Arabic technical terminology and maintain clarity.",
        "pt": "Focus on technical accuracy with appropriate Portuguese terminology.",
        "it": "Use precise Italian technical terminology with clarity.",
        "en": "Use precise technical English terminology and maintain accuracy.",
        # Default for other languages
        "default": "Prioritize technical precision and use appropriate specialized terminology."
    },
    "marketing": {
        "ja": "Use engaging Japanese with cultural nuance for marketing purposes.",
        "es": "Create compelling Spanish copy with emotional resonance for marketing.",
        "fr": "Use persuasive French with cultural appeal for marketing materials.",
        "de": "Craft engaging German content suitable for marketing contexts.",
        "zh": "Use persuasive Chinese with cultural relevance for marketing.",
        "ko": "Create impactful Korean marketing content with cultural awareness.",
        "ru": "Use persuasive Russian suited for marketing and promotional content.",
        "ar": "Craft engaging Arabic with cultural nuance for marketing purposes.",
        "pt": "Use persuasive Portuguese with cultural appeal for marketing.",
        "it": "Create compelling Italian copy with emotional appeal for marketing.",
        "en": "Create engaging, persuasive English copy for marketing purposes.",
        # Default for other languages
        "default": "Create compelling, persuasive content suitable for marketing purposes."
    },
    # Add more style categories as needed
    "default": {
        "default": "Translate accurately while maintaining the original tone and intent."
    }
}

def extract_json(response_text):
    """Extract JSON from markdown code blocks or raw text"""
    # Try to find JSON code blocks
    matches = re.findall(r'```(?:json)?\n(.*?)\n```', response_text, re.DOTALL)
    if matches:
        return matches[0]
    # If no code blocks, try to find first JSON structure
    match = re.search(r'{(.*?)}', response_text, re.DOTALL)
    if match:
        return f'{{{match.group(1)}}}'
    return response_text

def get_style_prompt(style, target_lang):
    """Get the appropriate style prompt for the target language"""
    style_dict = STYLE_PROMPTS.get(style, STYLE_PROMPTS["default"])
    return style_dict.get(target_lang, style_dict.get("default", ""))

def translate_text(text, model, source_lang, target_lang, style_prompt):
    """Translate text using Gemini API with JSON format"""
    if not text.strip():
        print(f"{Fore.YELLOW}⏭️ Skipping empty text block{Style.RESET_ALL}")
        return text

    batch_data = [{"id": 0, "text": text}]  # Using id 0 as we are translating single text at a time for PPTX
    
    # Get full language names for clarity in the prompt
    source_lang_name = LANGUAGE_MAPPINGS.get(source_lang, source_lang)
    target_lang_name = LANGUAGE_MAPPINGS.get(target_lang, target_lang)

    prompt = f"""
    Translate from {source_lang_name} to {target_lang_name}. Maintain formatting EXACTLY.
    {style_prompt}
    Return ONLY VALID JSON using this format:
    {{
        "translations": [
            {{"id": <original_id>, "translation": "<translated_text>"}}
        ]
    }}
    DO NOT USE MARKDOWN. Ensure proper JSON escaping.
    Input: {json.dumps(batch_data, ensure_ascii=False)}
    """

    print(f"{Fore.YELLOW}🔧 Translating text ({len(text)} chars): {text[:50]}...{Style.RESET_ALL}")

    try:
        response = model.generate_content(prompt)
        print(f"{Fore.GREEN}📥 Raw Response:{Style.RESET_ALL} {response.text[:150]}...")

        cleaned_response = extract_json(response.text)
        print(f"{Fore.BLUE}🧹 Cleaned Response:{Style.RESET_ALL} {cleaned_response[:150]}...")

        try:
            result = json.loads(cleaned_response)
        except json.JSONDecodeError as e:
            print(f"{Fore.RED}❌ JSON Error:{Style.RESET_ALL} {str(e)}")
            print(f"{Fore.MAGENTA}🧬 Response Fragment:{Style.RESET_ALL} {cleaned_response[:500]}{Style.RESET_ALL}")
            return text  # Return original text on JSON error

        if 'translations' not in result or not result['translations']:
            print(f"{Fore.RED}⚠️ Missing 'translations' in JSON response{Style.RESET_ALL}")
            return text

        translation_item = result['translations'][0]
        if translation_item.get('id') != 0:
            print(f"{Fore.RED}⚠️ Incorrect ID in JSON response, expecting id 0{Style.RESET_ALL}")
            return text

        translated_text = translation_item.get('translation', text)
        print(f"{Fore.GREEN}✅ Translated to: {translated_text[:50]}...{Style.RESET_ALL}")
        return translated_text

    except Exception as e:
        print(f"{Fore.RED}💥 Translation Error:{Style.RESET_ALL} {str(e)}")
        print(f"{Fore.MAGENTA}📋 Failed Text:{Style.RESET_ALL} {text[:200]}{Style.RESET_ALL}")
        return text  # Return original text on error


def process_shape(shape, model, source_lang, target_lang, style_prompt):
    """Shape processor with error handling"""
    try:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            print(f"{Fore.CYAN}⚙️ Processing group shape{Style.RESET_ALL}")
            for sub_shape in shape.shapes:
                process_shape(sub_shape, model, source_lang, target_lang, style_prompt)
            return

        if shape.has_text_frame:
            print(f"{Fore.CYAN}📝 Processing text shape: {shape.shape_type}{Style.RESET_ALL}")
            text_frame = shape.text_frame
            for paragraph in text_frame.paragraphs:
                for run in paragraph.runs:
                    if run.text.strip():
                        original = run.text
                        translated = translate_text(original, model, source_lang, target_lang, style_prompt)
                        run.text = translated

        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            print(f"{Fore.CYAN}📊 Processing table{Style.RESET_ALL}")
            table = shape.table
            for row_idx, row in enumerate(table.rows):
                for cell_idx, cell in enumerate(row.cells):
                    if cell.text_frame:
                        for paragraph in cell.text_frame.paragraphs:
                            for run in paragraph.runs:
                                if run.text.strip():
                                    original = run.text
                                    translated = translate_text(original, model, source_lang, target_lang, style_prompt)
                                    run.text = translated

    except Exception as e:
        print(f"{Fore.RED}❌ Error processing shape: {str(e)}{Style.RESET_ALL}")

def generate_output_filename(input_file, target_lang):
    """Generate output filename based on input file and target language"""
    input_path = Path(input_file)
    return str(input_path.with_name(f"{input_path.stem}_{target_lang}{input_path.suffix}"))

def process_presentation(input_file, output_file, model_name, source_lang, target_lang, style):
    """Main processor with file validation"""
    try:
        input_path = Path(input_file)
        if not output_file:
            output_file = generate_output_filename(input_file, target_lang)
            
        if not input_path.exists():
            print(f"{Fore.RED}❌ Input file not found: {input_file}{Style.RESET_ALL}")
            return False

        print(f"{Fore.CYAN}📖 Loading presentation: {input_file}{Style.RESET_ALL}")
        prs = Presentation(input_file)

        print(f"{Fore.CYAN}🔍 Found {len(prs.slides)} slides{Style.RESET_ALL}")

        # Initialize Gemini model outside the loop for efficiency
        genai.configure(api_key=os.getenv('GOOGLE_API_KEY'))
        model = genai.GenerativeModel(model_name)
        print(f"{Fore.GREEN}✨ Using Gemini model: {model_name}{Style.RESET_ALL}")
        
        # Get style prompt based on style and target language
        style_prompt = get_style_prompt(style, target_lang)
        print(f"{Fore.CYAN}🎯 Using translation style: {style} for {target_lang}{Style.RESET_ALL}")
        print(f"{Fore.CYAN}📝 Style prompt: {style_prompt}{Style.RESET_ALL}")

        for slide_idx, slide in enumerate(prs.slides):
            print(f"{Fore.CYAN}🖼️ Processing slide {slide_idx + 1}/{len(prs.slides)}{Style.RESET_ALL}")
            for shape in slide.shapes:
                process_shape(shape, model, source_lang, target_lang, style_prompt)

            if slide.has_notes_slide:
                print(f"{Fore.CYAN}📝 Processing notes slide {slide_idx + 1}{Style.RESET_ALL}")
                notes_slide = slide.notes_slide
                if notes_slide.notes_text_frame:
                    for paragraph in notes_slide.notes_text_frame.paragraphs:
                        for run in paragraph.runs:
                            if run.text.strip():
                                original = run.text
                                translated = translate_text(original, model, source_lang, target_lang, style_prompt)
                                run.text = translated

        print(f"{Fore.CYAN}💾 Saving to {output_file}{Style.RESET_ALL}")
        prs.save(output_file)
        print(f"{Fore.GREEN}✅ Translation completed successfully. Output saved to: {output_file}{Style.RESET_ALL}")
        return True

    except Exception as e:
        print(f"{Fore.RED}❌ Critical error: {str(e)}{Style.RESET_ALL}")
        print(f"{Fore.RED}Error details: {str(e)}{Style.RESET_ALL}")
        return False

def process_multi_language(input_file, model_name, source_lang, target_langs, style):
    """Process a presentation for multiple target languages"""
    success = True
    for target_lang in target_langs:
        output_file = generate_output_filename(input_file, target_lang)
        print(f"{Fore.CYAN}🌐 Processing translation to {LANGUAGE_MAPPINGS.get(target_lang, target_lang)}{Style.RESET_ALL}")
        if not process_presentation(input_file, output_file, model_name, source_lang, target_lang, style):
            success = False
            print(f"{Fore.RED}❌ Failed to process {target_lang} translation{Style.RESET_ALL}")
    return success

if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="PPTX Translator using Gemini API")
    parser.add_argument("input", help="Input PowerPoint file (.pptx)")
    parser.add_argument("--output", help="Output PowerPoint file (.pptx). If not specified, uses [input_filename]_[lang].pptx")
    parser.add_argument("--source-lang", default="en", help="Source language code (e.g., 'en', 'ja', 'es')")
    parser.add_argument("--target-lang", help="Target language code (e.g., 'en', 'ja', 'es')")
    parser.add_argument("--target-langs", nargs="+", help="Multiple target language codes (e.g., 'en ja es')")
    parser.add_argument("--gemini-model", default="gemini-1.5-flash", help="Gemini model name (e.g., gemini-1.5-flash, gemini-pro)")
    parser.add_argument("--style-prompt", default="business", choices=list(STYLE_PROMPTS.keys()), 
                         help="Translation style (e.g., 'business', 'casual', 'technical', 'marketing')")
    parser.add_argument("--list-styles", action="store_true", help="List available translation styles")
    parser.add_argument("--list-languages", action="store_true", help="List available language codes")

    args = parser.parse_args()
    
    # Handle informational arguments
    if args.list_styles:
        print(f"{Fore.CYAN}Available translation styles:{Style.RESET_ALL}")
        for style in STYLE_PROMPTS:
            print(f"  - {style}")
        sys.exit(0)
        
    if args.list_languages:
        print(f"{Fore.CYAN}Available language codes:{Style.RESET_ALL}")
        for code, name in LANGUAGE_MAPPINGS.items():
            print(f"  - {code}: {name}")
        sys.exit(0)
    
    # Validate arguments
    if not args.target_lang and not args.target_langs:
        parser.error("Either --target-lang or --target-langs must be specified")
    
    if args.target_lang and args.target_langs:
        parser.error("Cannot specify both --target-lang and --target-langs")
    
    # Process single or multiple languages
    if args.target_langs:
        if not process_multi_language(args.input, args.gemini_model, args.source_lang, args.target_langs, args.style_prompt):
            sys.exit(1)
    else:
        if not process_presentation(args.input, args.output, args.gemini_model, args.source_lang, args.target_lang, args.style_prompt):
            sys.exit(1)