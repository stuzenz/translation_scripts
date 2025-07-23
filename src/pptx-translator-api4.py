#!/usr/bin/env python3
"""
Modern PPTX Translator using the new Google Gen AI SDK (google-genai)

This version uses the latest google-genai package which replaces the deprecated 
google-generativeai package. It includes robust JSON parsing, better error handling,
and improved context awareness.

Dependencies: pip install google-genai>=1.7.0 python-pptx colorama
"""

import argparse
import sys
import os
import json
import re
import time
import logging
from pathlib import Path
from typing import Optional, Dict, Any, List
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from colorama import Fore, Style, init
from itertools import islice 

# New Google Gen AI SDK imports
try:
    from google import genai
    from google.genai import types
except ImportError:
    print(f"{Fore.RED}❌ Error: google-genai package not found.{Style.RESET_ALL}")
    print(f"{Fore.YELLOW}Please install with: pip install google-genai{Style.RESET_ALL}")
    sys.exit(1)

# Initialize colorama
init(autoreset=True)

# Configure logging
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(levelname)s - %(message)s'
)
logger = logging.getLogger(__name__)

# Language mapping for style adaptation
LANGUAGE_MAPPINGS = {
    "en": "English",
    "ja": "Japanese", 
    "es": "Spanish",
    "fr": "French",
    "de": "German",
    "zh": "Chinese",
    "ko": "Korean",
    "ru": "Russian",
    "ar": "Arabic",
    "pt": "Portuguese",
    "it": "Italian",
}

# Translation style prompts
STYLE_PROMPTS = {
    "business": {
        "ja": "Keep the translation polite, collaborative, and aligned with Japanese business norms.",
        "es": "Maintain a formal and professional tone suitable for Spanish business contexts.",
        "fr": "Use formal language with appropriate business terminology for French corporate environments.",
        "de": "Ensure formality with proper German business etiquette and terminology.",
        "zh": "Adapt to Chinese business communication norms with appropriate formal expressions.",
        "ko": "Follow Korean business hierarchy and formality conventions.",
        "ru": "Use formal Russian suitable for professional business communication.",
        "ar": "Employ formal Arabic with appropriate business terminology.",
        "pt": "Maintain formal Portuguese suitable for professional settings.",
        "it": "Use formal Italian appropriate for business communication.",
        "en": "Maintain professional business English with appropriate terminology.",
        "default": "Maintain a professional, formal tone appropriate for business contexts."
    },
    "casual": {
        "ja": "Use natural, conversational Japanese while maintaining appropriate politeness.",
        "es": "Use friendly, conversational Spanish appropriate for general audiences.",
        "fr": "Use everyday French with a warm, approachable tone.",
        "de": "Employ conversational German with a friendly tone.",
        "zh": "Use natural, everyday Chinese expressions with appropriate politeness.",
        "ko": "Use conversational Korean while maintaining appropriate respect levels.",
        "ru": "Use friendly, everyday Russian for general communication.",
        "ar": "Use conversational Arabic with appropriate informality.",
        "pt": "Use natural, everyday Portuguese expressions.",
        "it": "Use casual Italian with a friendly tone.",
        "en": "Use natural, conversational English with a friendly tone.",
        "default": "Use a casual, friendly tone while maintaining cultural appropriateness."
    },
    "technical": {
        "ja": "Prioritize technical accuracy in Japanese with appropriate terminology.",
        "es": "Use precise Spanish technical terminology and maintain factual accuracy.",
        "fr": "Employ specific French technical terms with precision and clarity.",
        "de": "Use accurate German technical terminology with precise language.",
        "zh": "Focus on technical precision in Chinese with appropriate specialized vocabulary.",
        "ko": "Use Korean technical terminology with precision and clarity.",
        "ru": "Employ specialized Russian technical vocabulary with accuracy.",
        "ar": "Use precise Arabic technical terminology and maintain clarity.",
        "pt": "Focus on technical accuracy with appropriate Portuguese terminology.",
        "it": "Use precise Italian technical terminology with clarity.",
        "en": "Use precise technical English terminology and maintain accuracy.",
        "default": "Prioritize technical precision and use appropriate specialized terminology."
    },
    "marketing": {
        "ja": "Use engaging Japanese with cultural nuance for marketing purposes.",
        "es": "Create compelling Spanish copy with emotional resonance for marketing.",
        "fr": "Use persuasive French with cultural appeal for marketing materials.",
        "de": "Craft engaging German content suitable for marketing contexts.",
        "zh": "Use persuasive Chinese with cultural relevance for marketing.",
        "ko": "Create impactful Korean marketing content with cultural awareness.",
        "ru": "Use persuasive Russian suited for marketing and promotional content.",
        "ar": "Craft engaging Arabic with cultural nuance for marketing purposes.",
        "pt": "Use persuasive Portuguese with cultural appeal for marketing.",
        "it": "Create compelling Italian copy with emotional appeal for marketing.",
        "en": "Create engaging, persuasive English copy for marketing purposes.",
        "default": "Create compelling, persuasive content suitable for marketing purposes."
    },
    "default": {
        "default": "Translate accurately while maintaining the original tone and intent."
    }
}

class JSONExtractor:
    """Advanced JSON extraction from LLM responses with multiple robust strategies"""
    
    @staticmethod
    def extract_json_robust(response_text: str) -> Optional[str]:
        """
        Extract JSON from various response formats with comprehensive fallback strategies
        """
        if not response_text or not response_text.strip():
            logger.warning("Empty response text provided to JSON extractor")
            return None
            
        # Strategy 1: Try to find JSON in markdown code blocks
        json_blocks = re.findall(r'```(?:json)?\s*(.*?)\s*```', response_text, re.DOTALL | re.IGNORECASE)
        for block in json_blocks:
            cleaned_block = block.strip()
            if JSONExtractor._is_valid_json(cleaned_block):
                logger.debug("Successfully extracted JSON from markdown code block")
                return cleaned_block
        
        # Strategy 2: Find JSON objects using balanced bracket matching
        json_obj = JSONExtractor._extract_balanced_json(response_text)
        if json_obj and JSONExtractor._is_valid_json(json_obj):
            logger.debug("Successfully extracted JSON using balanced bracket matching")
            return json_obj
            
        # Strategy 3: Try to find and fix common JSON issues
        potential_json = JSONExtractor._find_and_fix_json(response_text)
        if potential_json and JSONExtractor._is_valid_json(potential_json):
            logger.debug("Successfully extracted and fixed JSON structure")
            return potential_json
            
        # Strategy 4: Extract potential JSON and attempt repairs
        repaired_json = JSONExtractor._extract_and_repair_json(response_text)
        if repaired_json and JSONExtractor._is_valid_json(repaired_json):
            logger.debug("Successfully repaired and extracted JSON")
            return repaired_json
            
        # Strategy 5: Last resort - try the original simple regex
        match = re.search(r'\{.*\}', response_text, re.DOTALL)
        if match:
            candidate = match.group(0)
            if JSONExtractor._is_valid_json(candidate):
                logger.debug("Successfully extracted JSON using simple regex fallback")
                return candidate
                
        logger.error(f"Failed to extract valid JSON from response: {response_text[:200]}...")
        return None
    
    @staticmethod
    def _extract_balanced_json(text: str) -> Optional[str]:
        """Extract JSON using balanced bracket matching"""
        start_idx = text.find('{')
        if start_idx == -1:
            return None
            
        bracket_count = 0
        in_string = False
        escape_next = False
        
        for i, char in enumerate(text[start_idx:], start_idx):
            if escape_next:
                escape_next = False
                continue
                
            if char == '\\':
                escape_next = True
                continue
                
            if char == '"' and not escape_next:
                in_string = not in_string
                continue
                
            if not in_string:
                if char == '{':
                    bracket_count += 1
                elif char == '}':
                    bracket_count -= 1
                    if bracket_count == 0:
                        return text[start_idx:i+1]
        
        return None
    
    @staticmethod
    def _find_and_fix_json(text: str) -> Optional[str]:
        """Find JSON-like structures and attempt to fix common formatting issues"""
        # Look for patterns that look like our expected structure
        pattern = r'\{\s*["\']translations["\']\s*:\s*\[.*?\]\s*\}'
        match = re.search(pattern, text, re.DOTALL | re.IGNORECASE)
        
        if match:
            candidate = match.group(0)
            return JSONExtractor._apply_json_fixes(candidate)
            
        return None
    
    @staticmethod
    def _extract_and_repair_json(text: str) -> Optional[str]:
        """Extract potential JSON and apply comprehensive repairs"""
        # Find any structure that starts with { and contains "translations"
        pattern = r'\{[^{}]*"translations"[^{}]*\}'
        matches = re.findall(pattern, text, re.DOTALL | re.IGNORECASE)
        
        for match in matches:
            repaired = JSONExtractor._apply_json_fixes(match)
            if JSONExtractor._is_valid_json(repaired):
                return repaired
        
        # Try with broader patterns
        pattern = r'\{.*?"translations".*?\}'
        match = re.search(pattern, text, re.DOTALL | re.IGNORECASE)
        if match:
            return JSONExtractor._apply_json_fixes(match.group(0))
            
        return None
    
    @staticmethod
    def _apply_json_fixes(candidate: str) -> str:
        """Apply a comprehensive set of JSON fixes"""
        # Fix single quotes to double quotes
        candidate = re.sub(r"'([^']*)':", r'"\1":', candidate)
        candidate = re.sub(r":\s*'([^']*)'", r': "\1"', candidate)
        
        # Fix missing quotes around keys
        candidate = re.sub(r'(\w+):', r'"\1":', candidate)
        
        # Fix trailing commas
        candidate = re.sub(r',\s*}', '}', candidate)
        candidate = re.sub(r',\s*]', ']', candidate)
        
        # Ensure proper spacing
        candidate = re.sub(r'{\s*', '{ ', candidate)
        candidate = re.sub(r'\s*}', ' }', candidate)
        candidate = re.sub(r'\[\s*', '[ ', candidate)
        candidate = re.sub(r'\s*]', ' ]', candidate)
        
        # Fix missing closing brackets/braces (common truncation issue)
        open_braces = candidate.count('{')
        close_braces = candidate.count('}')
        open_brackets = candidate.count('[')
        close_brackets = candidate.count(']')
        
        # Add missing closing brackets/braces
        if open_braces > close_braces:
            candidate += '}' * (open_braces - close_braces)
        if open_brackets > close_brackets:
            candidate += ']' * (open_brackets - close_brackets)
        
        return candidate
    
    @staticmethod
    def _is_valid_json(text: str) -> bool:
        """Check if text is valid JSON"""
        if not text or not text.strip():
            return False
            
        try:
            json.loads(text.strip())
            return True
        except (json.JSONDecodeError, ValueError):
            return False

class ModernGeminiTranslator:
    """Modern Gemini translator using the new google-genai SDK with enhanced error handling"""
    
    def __init__(self, model_name: str = "gemini-2.0-flash", max_retries: int = 3):
        self.model_name = model_name
        self.max_retries = max_retries
        self.client = None
        self._initialize_client()
        
    def _initialize_client(self):
        """Initialize the Gemini client with proper configuration"""
        try:
            api_key = os.getenv('GOOGLE_API_KEY')
            if not api_key:
                raise ValueError("GOOGLE_API_KEY environment variable not set")
            
            # Create client with the new SDK
            self.client = genai.Client(api_key=api_key)
            
            logger.info(f"Successfully initialized Gemini client with model: {self.model_name}")
            
        except Exception as e:
            logger.error(f"Failed to initialize Gemini client: {e}")
            raise
    
    def translate_text(self, text: str, source_lang: str, target_lang: str, 
                      style_prompt: str, context: Optional[List[str]] = None) -> str:
        """
        Translate text with robust error handling and retry logic using the new SDK
        """
        if not text.strip():
            logger.debug("Skipping empty text block")
            return text

        # Build the translation prompt
        prompt = self._build_translation_prompt(text, source_lang, target_lang, style_prompt, context)
        
        for attempt in range(self.max_retries + 1):
            try:
                logger.debug(f"Translation attempt {attempt + 1}/{self.max_retries + 1} for text: {text[:50]}...")
                
                # Use the new SDK API
                config = types.GenerateContentConfig(
                    temperature=0.1,  # Low temperature for consistent translations
                    top_p=0.8,
                    top_k=40,
                    max_output_tokens=2048,
                )
                
                response = self.client.models.generate_content(
                    model=self.model_name,
                    contents=prompt,
                    config=config
                )
                
                if not response or not response.text:
                    logger.warning(f"Empty response on attempt {attempt + 1}")
                    if attempt < self.max_retries:
                        time.sleep(1)
                        continue
                    else:
                        return text
                
                logger.debug(f"Raw response: {response.text[:150]}...")
                
                # Extract and parse JSON
                json_text = JSONExtractor.extract_json_robust(response.text)
                if not json_text:
                    logger.warning(f"Failed to extract JSON on attempt {attempt + 1}")
                    if attempt < self.max_retries:
                        time.sleep(1)
                        continue
                    else:
                        return text
                
                # Parse the JSON
                try:
                    result = json.loads(json_text)
                except json.JSONDecodeError as e:
                    logger.warning(f"JSON decode error on attempt {attempt + 1}: {e}")
                    logger.debug(f"Failed JSON: {json_text[:200]}...")
                    if attempt < self.max_retries:
                        time.sleep(1)
                        continue
                    else:
                        return text
                
                # Validate and extract translation
                translation = self._extract_translation_from_result(result, text)
                if translation is not None:
                    logger.debug(f"Successfully translated: {translation[:50]}...")
                    return translation
                    
                logger.warning(f"Invalid translation result on attempt {attempt + 1}")
                if attempt < self.max_retries:
                    time.sleep(1)
                    continue
                    
            except Exception as e:
                logger.error(f"Translation error on attempt {attempt + 1}: {e}")
                if attempt < self.max_retries:
                    time.sleep(2 ** attempt)  # Exponential backoff
                    continue
        
        logger.error(f"All translation attempts failed for text: {text[:100]}...")
        return text  # Return original text as fallback
    
    def _build_translation_prompt(self, text: str, source_lang: str, target_lang: str, 
                                style_prompt: str, context: Optional[List[str]] = None) -> str:
        """Build a comprehensive translation prompt optimized for the new SDK"""
        
        source_lang_name = LANGUAGE_MAPPINGS.get(source_lang, source_lang)
        target_lang_name = LANGUAGE_MAPPINGS.get(target_lang, target_lang)
        
        # Build context section
        context_section = ""
        if context:
            # Limit context to prevent token overflow
            limited_context = context[:8]  # Reduced to 8 for better performance
            context_section = f"""
CONTEXT for translation accuracy:
{json.dumps(limited_context, ensure_ascii=False, indent=2)}

Use this context to maintain consistency and terminology accuracy.
"""
        
        # Build the prompt with explicit structure requirements
        prompt = f"""You are a professional translator specializing in {source_lang_name} to {target_lang_name} translation.

TRANSLATION TASK:
Translate the following text from {source_lang_name} to {target_lang_name}.

STYLE REQUIREMENTS: {style_prompt}

{context_section}

CRITICAL OUTPUT REQUIREMENTS:
1. Return ONLY valid JSON in this EXACT format (no markdown, no explanations):
{{"translations": [{{"id": 0, "translation": "your_translated_text_here"}}]}}

2. Maintain ALL original formatting (spaces, punctuation, line breaks)
3. Use proper JSON escaping for special characters (quotes, backslashes, etc.)
4. If text is already in {target_lang_name}, return it unchanged
5. Preserve meaning and tone of the original

INPUT TEXT TO TRANSLATE:
{json.dumps([{"id": 0, "text": text}], ensure_ascii=False)}

OUTPUT (valid JSON only):"""

        return prompt
    
    def _extract_translation_from_result(self, result: Dict[str, Any], original_text: str) -> Optional[str]:
        """Extract translation from the parsed JSON result with enhanced validation"""
        
        if not isinstance(result, dict):
            logger.warning(f"Result is not a dictionary: {type(result)}")
            return None
            
        if 'translations' not in result:
            logger.warning("Missing 'translations' key in result")
            return None
            
        translations = result['translations']
        if not isinstance(translations, list) or not translations:
            logger.warning("'translations' is not a non-empty list")
            return None
            
        translation_item = translations[0]
        if not isinstance(translation_item, dict):
            logger.warning("Translation item is not a dictionary")
            return None
            
        if translation_item.get('id') != 0:
            logger.warning(f"Unexpected ID in translation: {translation_item.get('id')}")
            return None
            
        translation = translation_item.get('translation')
        if translation is None:
            logger.warning("Missing 'translation' key in translation item")
            return None
            
        # Additional validation: ensure translation is not empty or just whitespace
        if not str(translation).strip():
            logger.warning("Translation is empty or whitespace")
            return original_text
            
        return translation

def get_style_prompt(style: str, target_lang: str) -> str:
    """Get the appropriate style prompt for the target language"""
    style_dict = STYLE_PROMPTS.get(style, STYLE_PROMPTS["default"])
    return style_dict.get(target_lang, style_dict.get("default", ""))

def collect_slide_context(slide) -> List[str]:
    """Collect all text from a slide to provide context for translation"""
    context_items = []
    
    def extract_text(shape):
        items = []
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for sub_shape in shape.shapes:
                items.extend(extract_text(sub_shape))
        elif shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if text:
                items.append(text)
        elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            table = shape.table
            for row in table.rows:
                for cell in row.cells:
                    if cell.text_frame and cell.text_frame.text.strip():
                        items.append(cell.text_frame.text.strip())
        return items
    
    # Get all text from all shapes in the slide
    for shape in slide.shapes:
        context_items.extend(extract_text(shape))
    
    # Add notes if available
    if hasattr(slide, 'has_notes_slide') and slide.has_notes_slide and slide.notes_slide.notes_text_frame:
        notes_text = slide.notes_slide.notes_text_frame.text.strip()
        if notes_text:
            context_items.append(f"[Slide Notes] {notes_text}")
            
    return context_items

def process_shape(shape, translator: ModernGeminiTranslator, source_lang: str, target_lang: str, 
                 style_prompt: str, context: Optional[List[str]] = None):
    """Shape processor with improved error handling and progress feedback"""
    try:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            print(f"{Fore.CYAN}⚙️ Processing group shape{Style.RESET_ALL}")
            for sub_shape in shape.shapes:
                process_shape(sub_shape, translator, source_lang, target_lang, style_prompt, context)
            return

        if shape.has_text_frame:
            print(f"{Fore.CYAN}📝 Processing text shape: {shape.shape_type}{Style.RESET_ALL}")
            text_frame = shape.text_frame
            for paragraph in text_frame.paragraphs:
                for run in paragraph.runs:
                    if run.text.strip():
                        original = run.text
                        translated = translator.translate_text(original, source_lang, target_lang, style_prompt, context)
                        run.text = translated
                        print(f"{Fore.GREEN}✅ '{original[:30]}...' → '{translated[:30]}...'{Style.RESET_ALL}")

        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            print(f"{Fore.CYAN}📊 Processing table{Style.RESET_ALL}")
            table = shape.table
            
            # First collect all text from the table for context
            table_context = []
            for row in table.rows:
                row_texts = []
                for cell in row.cells:
                    if cell.text_frame:
                        row_texts.append(cell.text_frame.text.strip())
                if row_texts:
                    table_context.append(" | ".join(row_texts))
            
            # Combine table-specific context with slide context
            enhanced_context = context.copy() if context else []
            if table_context:
                enhanced_context.append("[Table Content] " + "\n".join(table_context[:3]))  # Limit table context
            
            # Process each cell with the enhanced context
            for row_idx, row in enumerate(table.rows):
                for cell_idx, cell in enumerate(row.cells):
                    if cell.text_frame:
                        for paragraph in cell.text_frame.paragraphs:
                            for run in paragraph.runs:
                                if run.text.strip():
                                    original = run.text
                                    position_context = f"[Table Cell] Row {row_idx+1}, Column {cell_idx+1}"
                                    cell_context = enhanced_context + [position_context]
                                    translated = translator.translate_text(original, source_lang, target_lang, style_prompt, cell_context)
                                    run.text = translated
                                    print(f"{Fore.GREEN}✅ Table[{row_idx+1},{cell_idx+1}]: '{original[:20]}...' → '{translated[:20]}...'{Style.RESET_ALL}")

    except Exception as e:
        logger.error(f"Error processing shape: {e}")
        print(f"{Fore.RED}❌ Error processing shape: {str(e)}{Style.RESET_ALL}")

def generate_output_filename(input_file: str, target_lang: str) -> str:
    """Generate output filename based on input file and target language"""
    input_path = Path(input_file)
    return str(input_path.with_name(f"{input_path.stem}_{target_lang}{input_path.suffix}"))

def extract_presentation_title(prs) -> Optional[str]:
    """Extract the title of the presentation from the first slide if possible"""
    if not prs.slides or len(prs.slides) == 0:
        return None
        
    # Try to find a title in the first slide
    first_slide = prs.slides[0]
    for shape in first_slide.shapes:
        if shape.has_text_frame and shape.text_frame.text.strip():
            return shape.text_frame.text.strip()
    return None

def collect_global_context(prs, max_slides: int = 4) -> List[str]:
    """Build a focused pool of presentation‑wide context strings"""
    context = []

    title = extract_presentation_title(prs)
    if title:
        context.append(f"[Presentation Title] {title}")

    # Iterate safely with limited slides for context
    for i, slide in enumerate(islice(prs.slides, max_slides)):
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                context.append(f"[Slide {i+1} Title] {shape.text_frame.text.strip()}")
                break

    return context

def process_presentation(input_file: str, output_file: Optional[str], model_name: str, 
                        source_lang: str, target_lang: str, style: str) -> bool:
    """Main processor with enhanced error handling and modern SDK integration"""
    try:
        input_path = Path(input_file)
        if not output_file:
            output_file = generate_output_filename(input_file, target_lang)
            
        if not input_path.exists():
            print(f"{Fore.RED}❌ Input file not found: {input_file}{Style.RESET_ALL}")
            return False

        print(f"{Fore.CYAN}📖 Loading presentation: {input_file}{Style.RESET_ALL}")
        prs = Presentation(input_file)

        print(f"{Fore.CYAN}🔍 Found {len(prs.slides)} slides{Style.RESET_ALL}")

        # Initialize modern translator
        translator = ModernGeminiTranslator(model_name)
        print(f"{Fore.GREEN}✨ Using modern Gemini SDK with model: {model_name}{Style.RESET_ALL}")
        
        # Get style prompt based on style and target language
        style_prompt = get_style_prompt(style, target_lang)
        print(f"{Fore.CYAN}🎯 Using translation style: {style} for {target_lang}{Style.RESET_ALL}")
        print(f"{Fore.CYAN}📝 Style prompt: {style_prompt}{Style.RESET_ALL}")
        
        # Collect global context for the entire presentation
        global_context = collect_global_context(prs)
        print(f"{Fore.CYAN}🌍 Collected global context with {len(global_context)} items{Style.RESET_ALL}")

        # Process the presentation slide by slide
        for slide_idx, slide in enumerate(prs.slides):
            print(f"{Fore.CYAN}🖼️ Processing slide {slide_idx + 1}/{len(prs.slides)}{Style.RESET_ALL}")
            
            # Collect context from the current slide
            slide_context = collect_slide_context(slide)
            print(f"{Fore.CYAN}📄 Collected {len(slide_context)} context items from slide{Style.RESET_ALL}")
            
            # Combine global and slide-specific context (limit for performance)
            combined_context = global_context + [f"[Current Slide {slide_idx + 1}]"] + slide_context[:5]
            
            # Process all shapes with combined context
            for shape in slide.shapes:
                process_shape(shape, translator, source_lang, target_lang, style_prompt, combined_context)

            # Process notes with robust error handling
            try:
                if hasattr(slide, 'has_notes_slide') and slide.has_notes_slide:
                    print(f"{Fore.CYAN}📝 Processing notes slide {slide_idx + 1}{Style.RESET_ALL}")
                    
                    notes_slide = slide.notes_slide
                    if hasattr(notes_slide, 'notes_text_frame'):
                        text_frame = notes_slide.notes_text_frame
                        if hasattr(text_frame, 'paragraphs'):
                            for paragraph in text_frame.paragraphs:
                                if hasattr(paragraph, 'runs'):
                                    for run in paragraph.runs:
                                        if run.text.strip():
                                            original = run.text
                                            notes_context = combined_context + [f"[Notes for Slide {slide_idx + 1}]"]
                                            translated = translator.translate_text(original, source_lang, target_lang, style_prompt, notes_context)
                                            run.text = translated
                                            print(f"{Fore.GREEN}✅ Notes: '{original[:30]}...' → '{translated[:30]}...'{Style.RESET_ALL}")
                    
            except Exception as e:
                print(f"{Fore.YELLOW}⚠️ Error processing notes for slide {slide_idx + 1}: {str(e)}{Style.RESET_ALL}")
                continue

        print(f"{Fore.CYAN}💾 Saving to {output_file}{Style.RESET_ALL}")
        prs.save(output_file)
        print(f"{Fore.GREEN}✅ Translation completed successfully. Output saved to: {output_file}{Style.RESET_ALL}")
        return True

    except Exception as e:
        logger.error(f"Critical error in process_presentation: {e}")
        print(f"{Fore.RED}❌ Critical error: {str(e)}{Style.RESET_ALL}")
        return False

def process_multi_language(input_file: str, model_name: str, source_lang: str, 
                          target_langs: List[str], style: str) -> bool:
    """Process a presentation for multiple target languages"""
    success = True
    for target_lang in target_langs:
        output_file = generate_output_filename(input_file, target_lang)
        print(f"{Fore.CYAN}🌐 Processing translation to {LANGUAGE_MAPPINGS.get(target_lang, target_lang)}{Style.RESET_ALL}")
        if not process_presentation(input_file, output_file, model_name, source_lang, target_lang, style):
            success = False
            print(f"{Fore.RED}❌ Failed to process {target_lang} translation{Style.RESET_ALL}")
        else:
            print(f"{Fore.GREEN}✅ Successfully created {output_file}{Style.RESET_ALL}")
    return success

if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Modern PPTX Translator using Google Gen AI SDK")
    parser.add_argument("input", help="Input PowerPoint file (.pptx)")
    parser.add_argument("--output", help="Output PowerPoint file (.pptx)")
    parser.add_argument("--source-lang", default="en", help="Source language code")
    parser.add_argument("--target-lang", help="Target language code")
    parser.add_argument("--target-langs", nargs="+", help="Multiple target language codes")
    parser.add_argument("--gemini-model", default="gemini-2.0-flash", help="Gemini model name")
    parser.add_argument("--style-prompt", default="business", choices=list(STYLE_PROMPTS.keys()), 
                         help="Translation style")
    parser.add_argument("--max-retries", type=int, default=3, help="Maximum retries for failed translations")
    parser.add_argument("--debug", action="store_true", help="Enable detailed debug logging")
    parser.add_argument("--list-styles", action="store_true", help="List available translation styles")
    parser.add_argument("--list-languages", action="store_true", help="List available language codes")

    args = parser.parse_args()
    
    # Set logging level
    if args.debug:
        logging.getLogger().setLevel(logging.DEBUG)
        logger.debug("Debug logging enabled")
    
    # Handle informational arguments
    if args.list_styles:
        print(f"{Fore.CYAN}Available translation styles:{Style.RESET_ALL}")
        for style in STYLE_PROMPTS:
            print(f"  - {style}")
        sys.exit(0)
        
    if args.list_languages:
        print(f"{Fore.CYAN}Available language codes:{Style.RESET_ALL}")
        for code, name in LANGUAGE_MAPPINGS.items():
            print(f"  - {code}: {name}")
        sys.exit(0)
    
    # Validate arguments
    if not args.target_lang and not args.target_langs:
        parser.error("Either --target-lang or --target-langs must be specified")
    
    if args.target_lang and args.target_langs:
        parser.error("Cannot specify both --target-lang and --target-langs")

    # Show configuration summary
    print(f"{Fore.GREEN}=== Modern Translation Configuration ==={Style.RESET_ALL}")
    print(f"{Fore.CYAN}SDK Version:{Style.RESET_ALL} Google Gen AI SDK (modern)")
    print(f"{Fore.CYAN}Source Language:{Style.RESET_ALL} {args.source_lang}")
    if args.target_lang:
        print(f"{Fore.CYAN}Target Language:{Style.RESET_ALL} {args.target_lang}")
    else:
        print(f"{Fore.CYAN}Target Languages:{Style.RESET_ALL} {', '.join(args.target_langs)}")
    print(f"{Fore.CYAN}Translation Style:{Style.RESET_ALL} {args.style_prompt}")
    print(f"{Fore.CYAN}Gemini Model:{Style.RESET_ALL} {args.gemini_model}")
    print(f"{Fore.CYAN}Max Retries:{Style.RESET_ALL} {args.max_retries}")
    print(f"{Fore.GREEN}========================================={Style.RESET_ALL}")
    
    # Process single or multiple languages
    try:
        if args.target_langs:
            if not process_multi_language(args.input, args.gemini_model, args.source_lang, args.target_langs, args.style_prompt):
                sys.exit(1)
        else:
            if not process_presentation(args.input, args.output, args.gemini_model, args.source_lang, args.target_lang, args.style_prompt):
                sys.exit(1)
    except KeyboardInterrupt:
        print(f"\n{Fore.YELLOW}Translation interrupted by user{Style.RESET_ALL}")
        sys.exit(1)
    except Exception as e:
        logger.error(f"Unexpected error: {e}")
        print(f"{Fore.RED}❌ Unexpected error: {str(e)}{Style.RESET_ALL}")
        sys.exit(1)